#!/usr/bin/env python3
"""Fill Samsung's App Description template (v1.42) for Fliks.

Keeps the template's section order — UI Structure, Usage Scenario, Menu &
function description, Key Policy, How to change languages — and drops the
guidance/sample slides, as the template's first slide instructs.

    uv run --with python-pptx python store/tools/samsung-app-description.py
"""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Emu, Inches, Pt

SAMSUNG = Path(__file__).resolve().parents[1] / 'samsung'
SHOTS = SAMSUNG / 'screenshots'
OUT = SAMSUNG / 'app-description.pptx'

BG = RGBColor(0x1D, 0x23, 0x2A)
FG = RGBColor(0xFF, 0xFF, 0xFF)
DIM = RGBColor(0xB6, 0xBE, 0xC8)
ACCENT = RGBColor(0x7A, 0x3F, 0xF2)
TODO = RGBColor(0xFF, 0xC1, 0x07)

prs = Presentation()
prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
BLANK = prs.slide_layouts[6]
W, H = prs.slide_width, prs.slide_height


def slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(1, 0, 0, W, H)  # 1 = rectangle
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    bg.line.fill.background()
    bg.shadow.inherit = False
    return s


def textbox(s, x, y, w, h):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    return tf


def para(tf, text, size=14, color=FG, bold=False, first=False, space_after=6, level=0):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.level = level
    p.space_after = Pt(space_after)
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = 'Verdana'
    return p


def title(s, text, sub=None):
    tf = textbox(s, 0.6, 0.45, 12.1, 1.0)
    para(tf, text, size=30, bold=True, first=True, space_after=2)
    bar = s.shapes.add_shape(1, Inches(0.6), Inches(1.28), Inches(1.6), Pt(4))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()
    bar.shadow.inherit = False
    if sub:
        para(textbox(s, 0.6, 1.45, 12.1, 0.5), sub, size=13, color=DIM, first=True)


def bullets(s, items, x=0.6, y=2.0, w=12.1, h=5.0, size=14):
    tf = textbox(s, x, y, w, h)
    for i, item in enumerate(items):
        if isinstance(item, tuple):
            text, level = item
        else:
            text, level = item, 0
        para(
            tf,
            ('• ' if level == 0 else '– ') + text,
            size=size if level == 0 else size - 1,
            color=FG if level == 0 else DIM,
            first=(i == 0),
            level=level,
            space_after=7,
        )
    return tf


def table(s, rows, x=0.6, y=2.0, w=12.1, col_widths=None, size=12):
    n_rows, n_cols = len(rows), len(rows[0])
    shape = s.shapes.add_table(n_rows, n_cols, Inches(x), Inches(y), Inches(w), Inches(0.4 * n_rows))
    tbl = shape.table
    if col_widths:
        total = sum(col_widths)
        for i, cw in enumerate(col_widths):
            tbl.columns[i].width = Emu(int(Inches(w).inches * 914400 * cw / total))
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = tbl.cell(r, c)
            cell.text = ''
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0x2A, 0x31, 0x3A) if r == 0 else BG
            p = cell.text_frame.paragraphs[0]
            run = p.add_run()
            run.text = val
            run.font.size = Pt(size)
            run.font.bold = r == 0
            run.font.color.rgb = TODO if '<<' in val else FG
            run.font.name = 'Verdana'
    return tbl


# ---------------------------------------------------------------- slide 1: title
s = slide()
tf = textbox(s, 0.9, 2.6, 11.5, 2.2)
para(tf, 'Fliks', size=54, bold=True, first=True, space_after=4)
para(tf, 'App Description', size=28, color=DIM)
para(tf, 'Samsung Smart TV (Tizen) — package FliksMedia', size=15, color=DIM, space_after=2)
para(tf, 'CP name: Clément Delestre', size=15, color=DIM)

# ------------------------------------------------------- slide 2: revision history
s = slide()
title(s, 'Revision History')
table(
    s,
    [
        ['Version', 'Date', 'Description', 'Author'],
        ['1.0', '2026.08.03', 'First submission of Fliks 2.0.0', 'Clément Delestre'],
    ],
    col_widths=[1, 1.4, 5, 2],
)

# ---------------------------------------------------------------- slide 3: contents
s = slide()
title(s, 'Contents')
bullets(
    s,
    [
        'What the app is, and what it needs to run',
        'UI Structure',
        'Usage Scenario (server pairing and log-in)',
        'Menu & function description',
        'Key policy',
        'How to change languages',
    ],
    size=17,
)

# ------------------------------------------------- slide 4: what it is / test setup
s = slide()
title(s, 'What Fliks is', 'Read this first: the app cannot be tested without a server.')
bullets(
    s,
    [
        'Fliks plays the video files a user already owns. It is the TV client of a '
        'self-hosted media server the user runs at home.',
        'The app ships no content of its own and no catalogue: everything shown comes '
        'from the server the user pairs with. There is no purchase, no subscription '
        'and no in-app payment.',
        'On first launch the app therefore asks for a server address, then for an '
        'account on that server.',
        'A public test server is provided for validation — credentials on the next slide.',
        'Video is played with the Samsung AVPlay API (HLS, fMP4 and TS); the server '
        'converts the file on the fly when the TV cannot decode the original.',
    ],
    size=15,
)

# -------------------------------------------------------------- slide 5: UI structure
s = slide()
title(s, 'UI Structure', 'Depth navigation. Every screen is reached with the D-pad.')
bullets(
    s,
    [
        '1st depth — Server setup (first launch only): enter the server address',
        ('2nd depth — User picker: pick an account, or "Enter another user"', 1),
        ('3rd depth — Log in: username + password', 1),
        '1st depth — Home: Libraries, Continue watching, Recently added',
        ('2nd depth — Library (e.g. Movies): grid, Suggestions, Genres, Collections, Likes', 1),
        ('3rd depth — Media detail: synopsis, cast, file information, Play / Resume', 2),
        ('4th depth — Player: full-screen playback, audio and subtitle tracks, quality', 2),
        ('3rd depth — Series detail: seasons and episodes → 4th depth Player', 2),
        '1st depth — Sidebar: Home, Search, My profile, Playlists, History, Requests, '
        'Activity, Calendar',
        ('2nd depth — Search: results by title, person or genre → Media detail', 1),
        ('2nd depth — My profile / Settings: display, playback, languages, account', 1),
    ],
    size=13,
)

# ------------------------------------------------------------ slide 6: usage scenario
s = slide()
title(s, 'Usage Scenario', 'Log-in flow, as the template requires. No payment step exists.')
tf = bullets(
    s,
    [
        'Step 1 — Launch the app. On a first launch the "Server" screen is shown.',
        'Step 2 — Enter the test server address with the on-screen keyboard:',
        'Step 3 — The user picker appears ("Who\'s watching?"). Select the test account:',
        'Step 4 — Enter its password:',
        'Step 5 — Home is shown. Move focus with the D-pad, press ENTER to open a title.',
        'Step 6 — On the media detail screen, press ENTER on "Play" to start playback.',
        'Step 7 — During playback: LEFT/RIGHT seek, ENTER pauses and resumes, '
        'RETURN goes back to the detail screen.',
        'Step 8 — RETURN from Home closes the app (Fliks is not a game app, so no exit '
        'pop-up is shown); EXIT closes it as well and returns to broadcasting.',
    ],
    y=1.95,
    size=13,
)
tf2 = textbox(s, 1.1, 5.75, 11.0, 1.4)
para(tf2, 'Test account for validation', size=13, bold=True, color=DIM, first=True)
para(tf2, 'Server address:  fliks-demo.delestre.me', size=14)
para(tf2, 'User:  user-demo        Password:  user-demo', size=14)

# --------------------------------------------- slides 7-10: menu & function description
SCREENS = [
    (
        '1-home.jpg',
        'Home',
        [
            'Libraries the account can see, "Continue watching" with the resume '
            'position of every unfinished title, and "Recently added".',
            'Focus moves with UP/DOWN between rows and LEFT/RIGHT inside a row; '
            'ENTER opens the focused title.',
        ],
    ),
    (
        '4-library.jpg',
        'Library and navigation menu',
        [
            'The left menu reaches Home, Search, the libraries, Playlists, Downloads, '
            'History, Requests, Activity and Calendar.',
            'A library lists its titles as a poster grid, with sorting and the '
            'Suggestions, Genres, Collections and Likes tabs.',
        ],
    ),
    (
        '2-movie-detail.jpg',
        'Media detail',
        [
            'Poster, rating, year, runtime, genres and synopsis, plus the audio and '
            'subtitle tracks found in the file.',
            '"Resume" continues at the stored position, "From the beginning" restarts. '
            '"File information" shows the real codecs, bitrate and resolution.',
        ],
    ),
    (
        '3-playback.jpg',
        'Player',
        [
            'Full-screen playback through the Samsung AVPlay API. The bar shows the '
            'elapsed and total time and the buffered range.',
            'Controls: play/pause, skip back and forward, subtitle track, audio track, '
            'and playback settings (quality, speed).',
        ],
    ),
]
for img, name, notes in SCREENS:
    s = slide()
    title(s, f'Menu & function description — {name}')
    s.shapes.add_picture(str(SHOTS / img), Inches(0.6), Inches(1.95), width=Inches(7.6))
    tf = textbox(s, 8.5, 1.95, 4.2, 4.6)
    para(tf, name, size=18, bold=True, first=True, space_after=10)
    for n in notes:
        para(tf, '• ' + n, size=12, color=DIM, space_after=10)

# ---------------------------------------------------------------- slide 11: key policy
s = slide()
title(s, 'Key Policy', 'Every key below is handled by the app; the rest is left to the TV.')
table(
    s,
    [
        ['Button', 'Action', 'Remarks'],
        ['UP / DOWN', 'Move focus between rows, menu items and form fields', ''],
        ['LEFT / RIGHT', 'Move focus inside a row. During playback: seek backward / forward', ''],
        ['ENTER', 'Activate the focused item. On the player: pause and resume', ''],
        [
            'RETURN',
            'Close the topmost overlay, then go back one depth. At the top level: '
            'close the app (tizen.application.exit)',
            'Samsung Mandatory',
        ],
        ['EXIT', 'Close the app and return to broadcasting', 'Samsung Mandatory'],
        ['Ch. Up/Down', 'N/R', ''],
        ['Colour buttons', 'N/R', ''],
    ],
    col_widths=[1.6, 6.5, 2.2],
    size=12,
)

# ------------------------------------------------------------ slide 12: how to change languages
s = slide()
title(s, 'How to change languages')
table(
    s,
    [
        ['Items', 'Contents'],
        [
            'Supported languages',
            'English, French, Spanish, German, Italian and Portuguese.',
        ],
        [
            'Default behaviour',
            'The app follows the TV OSD language. With the TV set to one of the six '
            'languages, the whole UI is shown in it; any other setting falls back to English.',
        ],
        [
            'How to check',
            'Change the TV OSD language (Settings > General > System Manager > Language) '
            'and relaunch the app — or override it in the app itself: '
            'My profile > Display > Language.',
        ],
        [
            'Media titles and synopses',
            'Served by the user\'s own server in its metadata language; independent of the '
            'UI language above.',
        ],
    ],
    col_widths=[2.2, 9.9],
    size=12,
)

prs.save(str(OUT))
print(f'wrote {OUT} — {len(prs.slides.__iter__.__self__._sldIdLst)} slides')
